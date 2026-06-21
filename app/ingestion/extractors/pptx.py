import io

from pptx import Presentation

from app.ingestion.extractors.base import BaseExtractor, ExtractedChunk


def _extract_slide_title(slide) -> str | None:
    """Ambil judul slide dari placeholder title jika ada."""
    if slide.shapes.title and slide.shapes.title.has_text_frame:
        return slide.shapes.title.text.strip() or None
    return None


def _extract_slide_text(slide) -> str:
    """Gabungkan semua teks dari semua shape dalam slide."""
    texts: list[str] = []
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            line = " ".join(run.text for run in para.runs).strip()
            if line:
                texts.append(line)
    return "\n".join(texts)


class PPTXExtractor(BaseExtractor):
    """
    Ekstrak teks dari PPTX per slide.
    Setiap slide menjadi satu ExtractedChunk.
    Menyimpan slide_number dan slide_title untuk source tracing.
    """

    def extract(self, file_bytes: bytes, filename: str) -> list[ExtractedChunk]:
        prs = Presentation(io.BytesIO(file_bytes))
        chunks: list[ExtractedChunk] = []
        total_slides = len(prs.slides)

        for slide_num, slide in enumerate(prs.slides, start=1):
            title = _extract_slide_title(slide)
            text = _extract_slide_text(slide)

            if not text.strip():
                continue

            chunks.append(
                ExtractedChunk(
                    text=text,
                    slide_number=slide_num,
                    slide_title=title,
                    metadata={"total_slides": total_slides, "filename": filename},
                )
            )

        return chunks
